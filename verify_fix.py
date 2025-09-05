#!/usr/bin/env python3
"""
Verify that the fill fix worked correctly
"""

from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE

def verify_fill_fix():
    """Verify white text shapes have correct fill types"""
    print("=== VERIFYING FILL FIX ===\n")
    
    # Check the recreated presentation
    prs = Presentation("demo_recreated.pptx")
    
    for slide_idx, slide in enumerate(prs.slides):
        print(f"=== SLIDE {slide_idx + 1} ===")
        
        for shape_idx, shape in enumerate(slide.shapes):
            if hasattr(shape, 'text_frame') and shape.text_frame.text.strip():
                text_content = shape.text_frame.text.strip()
                
                # Look for the specific white text shapes we identified
                if "General Atomics" in text_content or "MDA-250083" in text_content or "250083 – MDA" in text_content:
                    print(f"\nFound target shape: {shape.name}")
                    print(f"  Text: '{text_content[:50]}...'")
                    
                    # Check font colors
                    has_white_font = False
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text.strip() and hasattr(run, 'font') and run.font.color:
                                try:
                                    if hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                                        rgb_str = str(run.font.color.rgb)
                                        if rgb_str.upper() in ['FFFFFF', 'FEFFFF', 'FDFEFF']:
                                            has_white_font = True
                                            print(f"    Font color: #{rgb_str} (WHITE)")
                                            break
                                except:
                                    pass
                    
                    # Check fill type
                    if hasattr(shape, 'fill'):
                        fill = shape.fill
                        print(f"    Fill type: {fill.type}")
                        
                        if fill.type == MSO_FILL_TYPE.GRADIENT:
                            print(f"    *** PROBLEM: GRADIENT FILL DETECTED ***")
                            try:
                                if hasattr(fill, 'fore_color') and fill.fore_color.rgb:
                                    print(f"      Gradient fore_color: #{str(fill.fore_color.rgb)}")
                                if hasattr(fill, 'back_color') and fill.back_color.rgb:
                                    print(f"      Gradient back_color: #{str(fill.back_color.rgb)}")
                            except:
                                pass
                        elif fill.type == MSO_FILL_TYPE.BACKGROUND:
                            print(f"    ✓ Correct: BACKGROUND fill (no explicit fill)")
                        elif fill.type == MSO_FILL_TYPE.SOLID:
                            try:
                                rgb_str = str(fill.fore_color.rgb)
                                print(f"    ✓ SOLID fill: #{rgb_str}")
                            except:
                                print(f"    ✓ SOLID fill (color access error)")
                        else:
                            print(f"    Fill type: {fill.type}")
                    
                    print("-" * 50)

if __name__ == "__main__":
    verify_fill_fix()