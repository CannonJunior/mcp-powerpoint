#!/usr/bin/env python3
"""
Debug script to investigate shadow effects on specific text
"""

from pptx import Presentation

def debug_shadow_text():
    """Debug shadow effects on Opportunity Vision and C2I2 text"""
    print("=== DEBUGGING SHADOW TEXT ISSUES ===\n")
    
    # Check both original and recreated presentations
    presentations = [
        ("ORIGINAL", "MDA-250083-BNB-20250904.v1.RFI.pptx"),
        ("RECREATED", "demo_recreated.pptx")
    ]
    
    target_texts = ["Opportunity Vision", "Command and Control", "C2I2"]
    
    for pres_name, pres_path in presentations:
        print(f"=== {pres_name} PRESENTATION ===")
        try:
            prs = Presentation(pres_path)
            
            for slide_idx, slide in enumerate(prs.slides):
                for shape_idx, shape in enumerate(slide.shapes):
                    if hasattr(shape, 'text_frame') and shape.text_frame.text.strip():
                        text_content = shape.text_frame.text.strip()
                        
                        # Check if this shape contains target text
                        contains_target = any(target in text_content for target in target_texts)
                        
                        if contains_target:
                            print(f"\n--- SLIDE {slide_idx + 1}, SHAPE {shape_idx + 1} ---")
                            print(f"Shape name: {shape.name}")
                            print(f"Text: '{text_content}'")
                            
                            # Check detailed text formatting
                            for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                                print(f"\n  Paragraph {para_idx + 1}: '{paragraph.text}'")
                                
                                for run_idx, run in enumerate(paragraph.runs):
                                    if run.text.strip():
                                        print(f"    Run {run_idx + 1}: '{run.text}'")
                                        
                                        if hasattr(run, 'font'):
                                            font = run.font
                                            
                                            # Font name and size
                                            print(f"      Font name: {getattr(font, 'name', 'None')}")
                                            if hasattr(font, 'size') and font.size:
                                                size_pt = font.size / 12700
                                                print(f"      Font size: {size_pt:.1f}pt")
                                            else:
                                                print(f"      Font size: None/Default")
                                            
                                            # Font color
                                            try:
                                                if hasattr(font, 'color') and font.color:
                                                    if hasattr(font.color, 'rgb') and font.color.rgb:
                                                        rgb_str = str(font.color.rgb)
                                                        print(f"      Font color: #{rgb_str}")
                                                    else:
                                                        print(f"      Font color: {font.color} (no RGB)")
                                                else:
                                                    print(f"      Font color: None/Default")
                                            except Exception as e:
                                                print(f"      Font color error: {e}")
                                            
                                            # Font shadow - DETAILED CHECK
                                            try:
                                                shadow_attr = getattr(font, 'shadow', None)
                                                print(f"      Font shadow attribute: {shadow_attr}")
                                                print(f"      Font shadow type: {type(shadow_attr)}")
                                                
                                                if shadow_attr is not None:
                                                    print(f"      *** SHADOW DETECTED ***")
                                                    
                                                    # Try to get shadow properties
                                                    try:
                                                        if hasattr(shadow_attr, 'enabled'):
                                                            print(f"        Shadow enabled: {shadow_attr.enabled}")
                                                        if hasattr(shadow_attr, 'visible'):
                                                            print(f"        Shadow visible: {shadow_attr.visible}")
                                                        if hasattr(shadow_attr, 'style'):
                                                            print(f"        Shadow style: {shadow_attr.style}")
                                                        if hasattr(shadow_attr, 'offset_x'):
                                                            print(f"        Shadow offset_x: {shadow_attr.offset_x}")
                                                        if hasattr(shadow_attr, 'offset_y'):
                                                            print(f"        Shadow offset_y: {shadow_attr.offset_y}")
                                                        if hasattr(shadow_attr, 'color'):
                                                            try:
                                                                if hasattr(shadow_attr.color, 'rgb'):
                                                                    shadow_rgb = str(shadow_attr.color.rgb)
                                                                    print(f"        Shadow color: #{shadow_rgb}")
                                                                else:
                                                                    print(f"        Shadow color: {shadow_attr.color}")
                                                            except:
                                                                print(f"        Shadow color: Could not access")
                                                    except Exception as se:
                                                        print(f"        Shadow properties error: {se}")
                                                else:
                                                    print(f"      No shadow")
                                                    
                                            except Exception as e:
                                                print(f"      Font shadow error: {e}")
                                            
                                            # Font effects
                                            try:
                                                print(f"      Font bold: {getattr(font, 'bold', 'None')}")
                                                print(f"      Font italic: {getattr(font, 'italic', 'None')}")
                                                print(f"      Font underline: {getattr(font, 'underline', 'None')}")
                                            except Exception as e:
                                                print(f"      Font effects error: {e}")
                            
                            print("-" * 60)
            
        except Exception as e:
            print(f"Error processing {pres_name}: {e}")
        
        print("\n" + "="*80 + "\n")

if __name__ == "__main__":
    debug_shadow_text()