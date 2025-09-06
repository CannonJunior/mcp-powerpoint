#!/usr/bin/env python3
"""
Test to check if line width is actually being set to 0
"""

from pptx import Presentation

def test_line_width():
    """Check line widths in recreated presentation"""
    print("=== CHECKING LINE WIDTHS ===")
    
    prs = Presentation("demo_recreated.pptx")
    
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            if hasattr(shape, 'text_frame') and shape.text_frame.text.strip():
                text_content = shape.text_frame.text.strip()
                
                # Focus on white font shapes
                has_white_font = False
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text.strip() and hasattr(run, 'font') and run.font.color:
                            try:
                                if hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                                    rgb_str = str(run.font.color.rgb)
                                    if rgb_str.upper() in ['FFFFFF', 'FEFFFF', 'FDFEFF']:
                                        has_white_font = True
                                        break
                            except:
                                pass
                
                if has_white_font:
                    print(f"\nShape {shape_idx + 1}: {shape.name}")
                    print(f"  Text: '{text_content[:30]}...'")
                    
                    if hasattr(shape, 'line'):
                        line = shape.line
                        try:
                            width = getattr(line, 'width', None)
                            print(f"  Line width: {width}")
                            
                            if width is not None:
                                width_pt = width / 12700
                                print(f"  Line width (pt): {width_pt:.3f}")
                                
                                if width_pt == 0:
                                    print(f"  ✓ Line width set to 0 - no border should be visible")
                                else:
                                    print(f"  ⚠ Line width > 0 - border may be visible")
                            else:
                                print(f"  Line width: None/Default")
                                
                        except Exception as e:
                            print(f"  Line width error: {e}")

if __name__ == "__main__":
    test_line_width()