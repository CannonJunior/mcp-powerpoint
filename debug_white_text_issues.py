#!/usr/bin/env python3
"""
Debug script to investigate white font text issues and blue border/shadow problems
"""

from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE, MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE_TYPE

def debug_white_text_issues():
    """Debug white font text issues with borders and shadows"""
    print("=== DEBUGGING WHITE TEXT ISSUES ===\n")
    
    # Check both original and recreated presentations
    presentations = [
        ("ORIGINAL", "MDA-250083-BNB-20250904.v1.RFI.pptx"),
        ("RECREATED", "demo_recreated.pptx")
    ]
    
    for pres_name, pres_path in presentations:
        print(f"=== {pres_name} PRESENTATION ===")
        try:
            prs = Presentation(pres_path)
            
            for slide_idx, slide in enumerate(prs.slides):
                print(f"\n--- SLIDE {slide_idx + 1} ---")
                
                for shape_idx, shape in enumerate(slide.shapes):
                    if hasattr(shape, 'text_frame') and shape.text_frame.text.strip():
                        text_content = shape.text_frame.text.strip()
                        has_white_font = False
                        has_black_font = False
                        
                        # Check if this shape has white or black text
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.text.strip() and hasattr(run, 'font') and run.font.color:
                                    try:
                                        if hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                                            rgb_str = str(run.font.color.rgb)
                                            if rgb_str.upper() in ['FFFFFF', 'FEFFFF', 'FDFEFF']:
                                                has_white_font = True
                                            elif rgb_str.upper() in ['000000', '000001', '010101']:
                                                has_black_font = True
                                    except:
                                        pass
                        
                        # Focus on shapes with white or black fonts
                        if has_white_font or has_black_font:
                            print(f"\nShape {shape_idx + 1}: {shape.name}")
                            print(f"  Text: '{text_content[:50]}...'")
                            print(f"  Has white font: {has_white_font}")
                            print(f"  Has black font: {has_black_font}")
                            
                            # Check detailed text formatting
                            for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                                for run_idx, run in enumerate(paragraph.runs):
                                    if run.text.strip():
                                        print(f"    Run {run_idx + 1}: '{run.text.strip()[:30]}...'")
                                        
                                        if hasattr(run, 'font'):
                                            font = run.font
                                            
                                            # Font color
                                            try:
                                                if hasattr(font, 'color') and font.color:
                                                    if hasattr(font.color, 'rgb') and font.color.rgb:
                                                        rgb_str = str(font.color.rgb)
                                                        print(f"      Font color: #{rgb_str}")
                                                    else:
                                                        print(f"      Font color: {font.color} (no RGB)")
                                                else:
                                                    print(f"      Font color: None/Default")
                                            except Exception as e:
                                                print(f"      Font color error: {e}")
                                            
                                            # Font shadow
                                            try:
                                                if hasattr(font, 'shadow') and font.shadow is not None:
                                                    print(f"      Font shadow: YES - {font.shadow}")
                                                else:
                                                    print(f"      Font shadow: None")
                                            except Exception as e:
                                                print(f"      Font shadow error: {e}")
                            
                            # Check shape border/line
                            if hasattr(shape, 'line'):
                                line = shape.line
                                print(f"    Shape border/line:")
                                try:
                                    # Line color
                                    if hasattr(line, 'color') and line.color:
                                        if hasattr(line.color, 'rgb') and line.color.rgb:
                                            line_rgb = str(line.color.rgb)
                                            print(f"      Border color: #{line_rgb}")
                                            if line_rgb.upper().startswith('3E7FCC') or line_rgb.upper().startswith('A4C1FF'):
                                                print(f"      *** PROBLEMATIC BLUE BORDER DETECTED: #{line_rgb} ***")
                                        else:
                                            print(f"      Border color: {line.color} (no RGB)")
                                    else:
                                        print(f"      Border color: None/Default")
                                        
                                    # Line width  
                                    if hasattr(line, 'width') and line.width:
                                        width_pt = line.width / 12700  # EMU to points
                                        print(f"      Border width: {width_pt:.1f}pt")
                                        if width_pt > 0:
                                            print(f"      *** BORDER PRESENT (width: {width_pt:.1f}pt) ***")
                                    else:
                                        print(f"      Border width: None/Default")
                                        
                                    # Line dash style
                                    if hasattr(line, 'dash_style'):
                                        print(f"      Border dash style: {line.dash_style}")
                                        
                                except Exception as e:
                                    print(f"      Border info error: {e}")
                            else:
                                print(f"    Shape border: No line attribute")
                            
                            # Check shape fill
                            if hasattr(shape, 'fill'):
                                fill = shape.fill
                                print(f"    Shape fill:")
                                try:
                                    print(f"      Fill type: {fill.type}")
                                    
                                    if fill.type == MSO_FILL_TYPE.SOLID:
                                        if hasattr(fill, 'fore_color') and fill.fore_color.rgb:
                                            fill_rgb = str(fill.fore_color.rgb)
                                            print(f"      Fill color: #{fill_rgb}")
                                    elif fill.type == MSO_FILL_TYPE.GRADIENT:
                                        print(f"      *** GRADIENT FILL DETECTED ***")
                                        if hasattr(fill, 'fore_color') and fill.fore_color.rgb:
                                            fore_rgb = str(fill.fore_color.rgb)
                                            print(f"        Gradient fore: #{fore_rgb}")
                                        if hasattr(fill, 'back_color') and fill.back_color.rgb:
                                            back_rgb = str(fill.back_color.rgb)
                                            print(f"        Gradient back: #{back_rgb}")
                                    elif fill.type == MSO_FILL_TYPE.BACKGROUND:
                                        print(f"      Background fill (no explicit fill)")
                                    
                                except Exception as e:
                                    print(f"      Fill info error: {e}")
                            
                            print("-" * 50)
            
        except Exception as e:
            print(f"Error processing {pres_name}: {e}")
        
        print("\n" + "="*80 + "\n")

if __name__ == "__main__":
    debug_white_text_issues()