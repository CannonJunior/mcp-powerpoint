#!/usr/bin/env python3
"""
Debug script to investigate fill color extraction issues
"""

from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE

def debug_fill_colors():
    """Debug fill color extraction"""
    prs = Presentation("MDA-250083-BNB-20250904.v1.RFI.pptx")
    
    print("=== DEBUGGING FILL COLOR EXTRACTION ===\n")
    
    slide = prs.slides[0]  # First slide
    
    for shape_idx in range(min(10, len(slide.shapes))):
        shape = slide.shapes[shape_idx]
        print(f"Shape {shape_idx + 1}: {shape.name} ({shape.shape_type})")
        
        if hasattr(shape, 'fill'):
            fill = shape.fill
            print(f"  Fill type: {fill.type}")
            
            try:
                if fill.type == MSO_FILL_TYPE.SOLID:
                    print("  SOLID fill detected")
                    if hasattr(fill, 'fore_color'):
                        print(f"    Has fore_color: {fill.fore_color}")
                        fore_color = fill.fore_color
                        if hasattr(fore_color, 'rgb'):
                            print(f"    RGB object: {fore_color.rgb}")
                            print(f"    RGB type: {type(fore_color.rgb)}")
                            try:
                                hex_val = f"#{fore_color.rgb:06X}"
                                print(f"    Hex color: {hex_val}")
                            except Exception as e:
                                print(f"    Hex conversion error: {e}")
                                print(f"    RGB str: {str(fore_color.rgb)}")
                        else:
                            print("    No RGB attribute")
                    else:
                        print("    No fore_color")
                        
                elif fill.type == MSO_FILL_TYPE.BACKGROUND:
                    print("  BACKGROUND fill detected")
                    
                else:
                    print(f"  Other fill type: {fill.type}")
                    
            except Exception as e:
                print(f"  Fill access error: {e}")
        else:
            print("  No fill attribute")
        
        print("-" * 50)

if __name__ == "__main__":
    debug_fill_colors()