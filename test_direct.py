#!/usr/bin/env python3
"""
Direct test of json_to_pptx function
"""

import sys
sys.path.append('.')

import json
from pptx import Presentation as PPTXPresentation
from pptx.util import Inches, Pt
from powerpoint_models import Presentation, Slide, Shape
from powerpoint_server import add_shape_to_slide, apply_text_formatting, apply_font_formatting

def test_direct():
    """Direct test of conversion"""
    print("=== DIRECT TEST OF JSON TO PPTX ===")
    
    # Read the JSON data
    with open('demo_presentation.json', 'r') as f:
        json_content = f.read()
    
    # Parse JSON
    data = json.loads(json_content)
    presentation_data = Presentation(**data)
    
    print("Creating presentation...")
    prs = PPTXPresentation()
    
    # Set slide dimensions
    prs.slide_width = presentation_data.slide_width
    prs.slide_height = presentation_data.slide_height
    
    for slide_data in presentation_data.slides:
        print(f"Processing slide...")
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        for shape_data in slide_data.shapes:
            if shape_data.text_frame and 'Financials' in shape_data.text_frame.text:
                print(f"Found Financials shape: {shape_data.name}")
                print(f"Text: {shape_data.text_frame.text[:50]}...")
                
                # Create the shape and apply formatting
                created_shape = add_shape_to_slide(slide, shape_data)
                if created_shape and hasattr(created_shape, 'text_frame'):
                    print("Applying text formatting...")
                    apply_text_formatting(created_shape.text_frame, shape_data.text_frame)
    
    # Save presentation
    prs.save("test_direct.pptx")
    print("Presentation saved as test_direct.pptx")

if __name__ == "__main__":
    test_direct()