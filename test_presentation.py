#!/usr/bin/env python3
"""
Create a simple test PowerPoint presentation for testing the web interface
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Create a presentation
prs = Presentation()

# Add a slide with title and content layout
slide_layout = prs.slide_layouts[1]  # Title and content layout
slide = prs.slides.add_slide(slide_layout)

# Set title
title = slide.shapes.title
title.text = "Project Status Overview"

# Add content
content = slide.placeholders[1]
content.text = "This presentation shows our quarterly results and upcoming goals."

# Add a text box with custom content
left = Inches(1)
top = Inches(3)
width = Inches(8)
height = Inches(1)

textbox = slide.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame
text_frame.text = "Key achievements this quarter include improved efficiency and customer satisfaction."

# Save the presentation
prs.save('/home/junior/src/mcp-powerpoint/test_simple.pptx')
print("Created test_simple.pptx")