#!/usr/bin/env python3
"""
Debug script to investigate text formatting and border issues
"""

from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE, MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE_TYPE

def debug_text_formatting():
    """Debug text formatting and border issues"""
    print("=== DEBUGGING TEXT FORMATTING AND BORDERS ===\n")
    
    # Check both original and recreated presentations
    presentations = [
        ("ORIGINAL", "MDA-250083-BNB-20250904.v1.RFI.pptx"),
        ("RECREATED", "demo_recreated.pptx")
    ]
    
    for pres_name, pres_path in presentations:
        print(f"=== {pres_name} PRESENTATION ===")
        try:
            prs = Presentation(pres_path)
            
            for slide_idx, slide in enumerate(prs.slides):
                print(f"\n--- SLIDE {slide_idx + 1} ---")
                
                for shape_idx, shape in enumerate(slide.shapes):
                    if hasattr(shape, 'text_frame') and shape.text_frame.text.strip():
                        text_content = shape.text_frame.text.strip()
                        
                        # Focus on specific problematic text
                        if any(keyword in text_content for keyword in ["General Atomics", "MDA-250083", "250083", "Financials"]):
                            print(f"\nShape {shape_idx + 1}: {shape.name}")
                            print(f"  Text: '{text_content[:50]}...'")
                            
                            # Check text formatting in detail
                            for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                                for run_idx, run in enumerate(paragraph.runs):
                                    if run.text.strip():
                                        print(f"    Run {run_idx + 1}: '{run.text.strip()}'")
                                        
                                        if hasattr(run, 'font'):
                                            font = run.font
                                            
                                            # Font color
                                            try:
                                                if hasattr(font, 'color') and font.color:
                                                    if hasattr(font.color, 'rgb') and font.color.rgb:
                                                        rgb_str = str(font.color.rgb)
                                                        print(f"      Font color: #{rgb_str}")
                                                    else:
                                                        print(f"      Font color: {font.color} (no RGB)")
                                                else:
                                                    print(f"      Font color: None/Default")
                                            except Exception as e:
                                                print(f"      Font color error: {e}")
                                            
                                            # Font shadow
                                            try:
                                                if hasattr(font, 'shadow') and font.shadow is not None:
                                                    print(f"      Font shadow: {font.shadow}")
                                                else:
                                                    print(f"      Font shadow: None")
                                            except Exception as e:
                                                print(f"      Font shadow error: {e}")
                            
                            # Check shape border/line
                            if hasattr(shape, 'line'):
                                line = shape.line
                                print(f"    Shape line/border:")
                                try:
                                    if hasattr(line, 'color') and line.color:
                                        if hasattr(line.color, 'rgb') and line.color.rgb:
                                            line_rgb = str(line.color.rgb)
                                            print(f"      Line color: #{line_rgb}")
                                        else:
                                            print(f"      Line color: {line.color} (no RGB)")
                                    else:
                                        print(f"      Line color: None/Default")
                                        
                                    if hasattr(line, 'width') and line.width:
                                        width_pt = line.width / 12700  # EMU to points
                                        print(f"      Line width: {width_pt:.1f}pt")
                                    else:
                                        print(f"      Line width: None/Default")
                                        
                                except Exception as e:
                                    print(f"      Line info error: {e}")
                            else:
                                print(f"    Shape line: No line attribute")
                            
                            print("-" * 40)
            
        except Exception as e:
            print(f"Error processing {pres_name}: {e}")
        
        print("\n" + "="*60 + "\n")

if __name__ == "__main__":
    debug_text_formatting()