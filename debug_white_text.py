#!/usr/bin/env python3
"""
Debug script to investigate white font text with incorrect gradient fills
"""

from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE

def debug_white_text():
    """Debug white text with gradient fill issues"""
    prs = Presentation("MDA-250083-BNB-20250904.v1.RFI.pptx")
    
    print("=== DEBUGGING WHITE TEXT GRADIENT FILLS ===\n")
    
    for slide_idx, slide in enumerate(prs.slides):
        print(f"=== SLIDE {slide_idx + 1} ===")
        
        for shape_idx, shape in enumerate(slide.shapes):
            # Focus on text shapes
            if hasattr(shape, 'text_frame') and shape.text_frame.text.strip():
                print(f"\nShape {shape_idx + 1}: {shape.name}")
                print(f"  Shape type: {shape.shape_type}")
                print(f"  Text: '{shape.text_frame.text.strip()}'")
                
                # Check text formatting
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text.strip():
                            print(f"    Run text: '{run.text}'")
                            if hasattr(run, 'font'):
                                font = run.font
                                # Check font color
                                if hasattr(font, 'color') and font.color:
                                    try:
                                        if hasattr(font.color, 'rgb') and font.color.rgb:
                                            rgb_str = str(font.color.rgb)
                                            print(f"      Font color: #{rgb_str}")
                                            
                                            # Check if this is white or near-white
                                            if rgb_str.upper() in ['FFFFFF', 'FEFFFF', 'FDFEFF']:
                                                print(f"      *** WHITE FONT DETECTED ***")
                                                
                                                # Check the shape fill
                                                if hasattr(shape, 'fill'):
                                                    fill = shape.fill
                                                    print(f"      Shape fill type: {fill.type}")
                                                    
                                                    if fill.type == MSO_FILL_TYPE.GRADIENT:
                                                        print(f"      *** GRADIENT FILL ON WHITE TEXT SHAPE ***")
                                                        try:
                                                            if hasattr(fill, 'fore_color') and fill.fore_color:
                                                                if hasattr(fill.fore_color, 'rgb'):
                                                                    fore_rgb = str(fill.fore_color.rgb)
                                                                    print(f"        Gradient fore_color: #{fore_rgb}")
                                                            if hasattr(fill, 'back_color') and fill.back_color:
                                                                if hasattr(fill.back_color, 'rgb'):
                                                                    back_rgb = str(fill.back_color.rgb)
                                                                    print(f"        Gradient back_color: #{back_rgb}")
                                                        except Exception as e:
                                                            print(f"        Gradient color extraction error: {e}")
                                                    
                                                    elif fill.type == MSO_FILL_TYPE.SOLID:
                                                        print(f"      Shape has SOLID fill")
                                                        try:
                                                            if hasattr(fill, 'fore_color') and fill.fore_color:
                                                                if hasattr(fill.fore_color, 'rgb'):
                                                                    fill_rgb = str(fill.fore_color.rgb)
                                                                    print(f"        Solid fill color: #{fill_rgb}")
                                                        except Exception as e:
                                                            print(f"        Solid color extraction error: {e}")
                                                            
                                                    elif fill.type == MSO_FILL_TYPE.BACKGROUND:
                                                        print(f"      Shape has BACKGROUND fill (no explicit fill)")
                                                        
                                    except Exception as e:
                                        print(f"      Font color error: {e}")
                
                print("-" * 40)
        
        print("\n" + "="*60 + "\n")

if __name__ == "__main__":
    debug_white_text()