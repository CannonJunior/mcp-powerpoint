#!/usr/bin/env python3
"""
Debug script specifically for Financials text formatting
"""

import json
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

def hex_to_rgb(hex_color: str):
    """Convert hex color string to RGBColor"""
    if not hex_color or not hex_color.startswith('#'):
        return None
    try:
        rgb_int = int(hex_color[1:], 16)
        return RGBColor((rgb_int >> 16) & 0xFF, (rgb_int >> 8) & 0xFF, rgb_int & 0xFF)
    except:
        return None

def debug_financials():
    """Debug Financials text specifically"""
    print("=== DEBUGGING FINANCIALS TEXT ===\n")
    
    # Check JSON data
    with open('demo_presentation.json', 'r') as f:
        data = json.load(f)
    
    print("1. CHECKING JSON EXTRACTION:")
    for slide_idx, slide in enumerate(data['slides']):
        for shape in slide['shapes']:
            if shape.get('text_frame') and 'Financials' in shape['text_frame'].get('text', ''):
                print(f"   Found Financials shape in slide {slide_idx + 1}")
                print(f"   Shape name: {shape['name']}")
                print(f"   Text: {shape['text_frame']['text'][:50]}...")
                
                for para_idx, paragraph in enumerate(shape['text_frame']['paragraphs']):
                    print(f"   Paragraph {para_idx + 1}: {paragraph['text']}")
                    for run_idx, run in enumerate(paragraph['runs']):
                        font = run.get('font', {})
                        print(f"     Run {run_idx + 1}: '{run['text']}' -> Color: {font.get('color_rgb')}, Size: {font.get('size')}")
    
    print("\n2. TESTING HEX CONVERSION:")
    test_color = "#000000"
    rgb_color = hex_to_rgb(test_color)
    print(f"   {test_color} -> {rgb_color}")
    
    print("\n3. CHECKING DEMO RECREATED PRESENTATION:")
    prs = Presentation("demo_recreated.pptx")
    
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            if hasattr(shape, 'text_frame') and 'Financials' in shape.text_frame.text:
                print(f"   Found Financials shape in slide {slide_idx + 1}")
                print(f"   Shape name: {shape.name}")
                
                for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                    for run_idx, run in enumerate(paragraph.runs):
                        if run.text.strip():
                            print(f"     Run {run_idx + 1}: '{run.text.strip()}'")
                            font = run.font
                            
                            # Font size
                            if hasattr(font, 'size') and font.size:
                                size_pt = font.size / 12700  # EMU to points
                                print(f"       Font size: {size_pt:.1f}pt")
                            else:
                                print(f"       Font size: None/Default")
                            
                            # Font color
                            if hasattr(font, 'color'):
                                try:
                                    if hasattr(font.color, 'rgb') and font.color.rgb:
                                        rgb_str = str(font.color.rgb)
                                        print(f"       Font color: #{rgb_str}")
                                    else:
                                        print(f"       Font color: {font.color} (no RGB access)")
                                except Exception as e:
                                    print(f"       Font color error: {e}")
    
    print("\n4. CHECKING DIRECT TEST RESULT:")
    try:
        prs_direct = Presentation("test_direct.pptx")
        
        for slide_idx, slide in enumerate(prs_direct.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                if hasattr(shape, 'text_frame') and 'Financials' in shape.text_frame.text:
                    print(f"   Found Financials shape in slide {slide_idx + 1}")
                    print(f"   Shape name: {shape.name}")
                    
                    for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                        for run_idx, run in enumerate(paragraph.runs):
                            if run.text.strip():
                                print(f"     Run {run_idx + 1}: '{run.text.strip()}'")
                                font = run.font
                                
                                # Font size
                                if hasattr(font, 'size') and font.size:
                                    size_pt = font.size / 12700  # EMU to points
                                    print(f"       Font size: {size_pt:.1f}pt")
                                else:
                                    print(f"       Font size: None/Default")
                                
                                # Font color
                                if hasattr(font, 'color'):
                                    try:
                                        if hasattr(font.color, 'rgb') and font.color.rgb:
                                            rgb_str = str(font.color.rgb)
                                            print(f"       Font color: #{rgb_str}")
                                        else:
                                            print(f"       Font color: {font.color} (no RGB access)")
                                    except Exception as e:
                                        print(f"       Font color error: {e}")
    except Exception as e:
        print(f"   Error reading test_direct.pptx: {e}")

if __name__ == "__main__":
    debug_financials()